#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 업데이트 스크립트
- 기존 PPT에 새로운 분석 결과 슬라이드 추가
- 이미지 삽입 및 텍스트 추가
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 경로 설정
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
RESULT_DIR = os.path.join(BASE_DIR, 'result_v3')
PPTX_PATH = os.path.join(RESULT_DIR, '산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v2.pptx')

def add_title_text(slide, text, left, top, width, height, font_size=24, bold=True, color=(0, 0, 0)):
    """제목 텍스트 추가"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return shape

def add_body_text(slide, text, left, top, width, height, font_size=12, color=(50, 50, 50)):
    """본문 텍스트 추가"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*color)
    return shape

def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """이미지 안전하게 추가"""
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top))
        return True
    else:
        print(f"  [경고] 이미지 없음: {img_path}")
        return False

def main():
    print("=" * 60)
    print("PPT 업데이트 시작")
    print("=" * 60)
    
    # PPT 로드
    prs = Presentation(PPTX_PATH)
    print(f"기존 슬라이드 수: {len(prs.slides)}")
    
    # 슬라이드 레이아웃 (빈 슬라이드)
    blank_layout = prs.slide_layouts[0]  # DEFAULT 레이아웃
    
    # =========================================================================
    # 새 슬라이드 1: 연도별 토지유형 변화량 분석
    # =========================================================================
    print("\n[1] 연도별 토지유형 변화량 분석 슬라이드 추가...")
    slide1 = prs.slides.add_slide(blank_layout)
    
    # 제목
    add_title_text(slide1, "연도별 토지유형 변화량 분석 (2017-2025)", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    
    # 부제목
    add_body_text(slide1, "임야/농경지는 감소, 대지/공장용지는 증가 추세 확인", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    # 이미지 추가
    img_path = os.path.join(RESULT_DIR, '05_yearly_landtype_change.png')
    add_image_safe(slide1, img_path, 0.3, 1.4, width=12.5, height=5.8)
    
    # 핵심 수치 텍스트
    summary_text = "▶ 2017→2025 변화: 임야 -0.84%, 농경지 -3.19%, 대지 +12.67%, 공장용지 +22.74%"
    add_body_text(slide1, summary_text, 0.5, 7.1, 12, 0.3, font_size=11, color=(200, 50, 50))
    
    # =========================================================================
    # 새 슬라이드 2: PCA 분류 근거
    # =========================================================================
    print("[2] PCA 분류 근거 슬라이드 추가...")
    slide2 = prs.slides.add_slide(blank_layout)
    
    add_title_text(slide2, "PCA 기반 지역 분류 근거", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    add_body_text(slide2, "주성분 분석으로 14개 시군구의 토지이용 특성을 2차원으로 시각화", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    img_path = os.path.join(RESULT_DIR, '02_PCA_classification_basis.png')
    add_image_safe(slide2, img_path, 0.3, 1.4, width=12.5, height=5.0)
    
    # 분석 결과 텍스트
    pca_text = """▶ PC1 (90.1%): 도시/산업화 정도 - 공장용지/대지 비율이 높을수록 양(+)의 방향
▶ PC2 (6.6%): 농경지 비율 - 농경지가 높을수록 양(+)의 방향  
▶ 누적 분산 설명력: 96.7% → 2개 주성분으로 토지이용 특성 대부분 설명"""
    add_body_text(slide2, pca_text, 0.5, 6.5, 12, 0.9, font_size=11, color=(50, 50, 50))
    
    # =========================================================================
    # 새 슬라이드 3: K-Means 실루엣 분석
    # =========================================================================
    print("[3] K-Means 실루엣 분석 슬라이드 추가...")
    slide3 = prs.slides.add_slide(blank_layout)
    
    add_title_text(slide3, "K-Means 군집화 타당성 검증", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    add_body_text(slide3, "Silhouette Score 0.552로 3개 군집 분류의 적절성 확인", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    img_path = os.path.join(RESULT_DIR, '03_silhouette_analysis.png')
    add_image_safe(slide3, img_path, 0.3, 1.4, width=12.5, height=4.5)
    
    silhouette_text = """▶ Elbow Method: k=3에서 최적 지점 확인
▶ Silhouette Score: k=3에서 0.552 (0.4 이상이면 군집 분리 적절)
▶ 군집별 Silhouette 분포: 모든 군집이 양(+)의 값으로 적절한 분류"""
    add_body_text(slide3, silhouette_text, 0.5, 6.0, 12, 0.9, font_size=11, color=(50, 50, 50))
    
    # =========================================================================
    # 새 슬라이드 4: 청주 4개구 비교 분석
    # =========================================================================
    print("[4] 청주 4개구 비교 분석 슬라이드 추가...")
    slide4 = prs.slides.add_slide(blank_layout)
    
    add_title_text(slide4, "청주시 4개구 토지이용 비교 분석", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    add_body_text(slide4, "같은 청주시 내에서도 구별로 상이한 토지이용 특성 확인", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    img_path = os.path.join(RESULT_DIR, '09_cheongju_4gu_comparison.png')
    add_image_safe(slide4, img_path, 0.3, 1.4, width=12.5, height=5.5)
    
    cheongju_text = """▶ 흥덕구/청원구: 도시/산업형 - 공장용지 비율 높음, 산업화 진행
▶ 서원구: 균형형 - 도시/농촌 기능 혼재
▶ 상당구: 농업/산림형 - 임야 76.7%로 산림 중심"""
    add_body_text(slide4, cheongju_text, 0.5, 7.0, 12, 0.5, font_size=11, color=(50, 50, 50))
    
    # =========================================================================
    # 새 슬라이드 5: 인구 상관관계 분석
    # =========================================================================
    print("[5] 인구 상관관계 분석 슬라이드 추가...")
    slide5 = prs.slides.add_slide(blank_layout)
    
    add_title_text(slide5, "인구 변화와 토지이용 상관관계", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    add_body_text(slide5, "인구 증가 지역에서 대지 면적 증가, 농경지 감소 통계적 유의", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    img_path = os.path.join(RESULT_DIR, '10_correlation_population.png')
    add_image_safe(slide5, img_path, 0.3, 1.4, width=12.5, height=5.0)
    
    pop_text = """▶ 인구↔대지: r=0.602, p=0.023 (유의) → 인구 증가 지역에서 대지 면적 증가
▶ 인구↔농경지: r=-0.668, p=0.009 (유의) → 인구 증가 지역에서 농경지 감소
▶ 인구↔공장용지: r=-0.122, p=0.678 (비유의) → 공장용지는 인구와 직접적 연관 없음"""
    add_body_text(slide5, pop_text, 0.5, 6.5, 12, 0.9, font_size=11, color=(50, 50, 50))
    
    # =========================================================================
    # 새 슬라이드 6: 도로 상관관계 분석
    # =========================================================================
    print("[6] 도로 상관관계 분석 슬라이드 추가...")
    slide6 = prs.slides.add_slide(blank_layout)
    
    add_title_text(slide6, "도로면적과 토지이용 상관관계", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    add_body_text(slide6, "도로 인프라 확충과 인구 증가 간 양의 상관관계 확인", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    img_path = os.path.join(RESULT_DIR, '11_correlation_road.png')
    add_image_safe(slide6, img_path, 0.3, 1.4, width=12.5, height=5.0)
    
    road_text = """▶ 도로↔인구: r=0.599, p=0.024 (유의) → 도로 인프라 확충 지역에서 인구 증가
▶ 도로↔대지: r=0.528, p=0.052 (경계) → 도로 확충과 대지 증가 연관성 있음
▶ 도로↔공장용지: r=0.476, p=0.086 (비유의) → 공장용지는 산업단지 입지 특성 반영"""
    add_body_text(slide6, road_text, 0.5, 6.5, 12, 0.9, font_size=11, color=(50, 50, 50))
    
    # =========================================================================
    # 새 슬라이드 7: 공장용지 성장 히트맵
    # =========================================================================
    print("[7] 공장용지 성장 히트맵 슬라이드 추가...")
    slide7 = prs.slides.add_slide(blank_layout)
    
    add_title_text(slide7, "시군구별 공장용지 성장지수 (2017=100)", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    add_body_text(slide7, "영동군, 청원구, 음성군 순으로 공장용지 성장률 높음", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    img_path = os.path.join(RESULT_DIR, '12_heatmap_factory_growth.png')
    add_image_safe(slide7, img_path, 0.5, 1.4, width=12.0, height=5.5)
    
    factory_text = """▶ 영동군: +47.95% (1위) - 신규 산업단지 조성 영향
▶ 청주 청원구: +34.05% (2위), 음성군: +25.62% (3위) - 기존 산업벨트 확장
▶ 농업/산림형 지역(보은, 단양)도 공장용지 소폭 증가"""
    add_body_text(slide7, factory_text, 0.5, 7.0, 12, 0.5, font_size=11, color=(50, 50, 50))
    
    # =========================================================================
    # 새 슬라이드 8: 유형별 토지 변화 추이
    # =========================================================================
    print("[8] 유형별 토지 변화 추이 슬라이드 추가...")
    slide8 = prs.slides.add_slide(blank_layout)
    
    add_title_text(slide8, "군집별 토지이용 비율 변화 추이", 0.5, 0.3, 12, 0.6, font_size=28, bold=True)
    add_body_text(slide8, "도시/산업형 지역에서 공장용지 증가 추세 뚜렷", 0.5, 0.85, 12, 0.4, font_size=14, color=(100, 100, 100))
    
    img_path = os.path.join(RESULT_DIR, '06_landuse_trend_by_cluster.png')
    add_image_safe(slide8, img_path, 0.3, 1.4, width=12.5, height=5.5)
    
    trend_text = """▶ 도시/산업형: 공장용지 지속 증가, 임야/농경지 감소 뚜렷
▶ 농업/산림형: 변화 폭 미미, 기존 토지이용 구조 유지
▶ 균형형: 대지 증가 추세, 도시화 진행 중"""
    add_body_text(slide8, trend_text, 0.5, 7.0, 12, 0.5, font_size=11, color=(50, 50, 50))
    
    # =========================================================================
    # 저장
    # =========================================================================
    output_path = PPTX_PATH
    prs.save(output_path)
    
    print("\n" + "=" * 60)
    print("PPT 업데이트 완료!")
    print(f"최종 슬라이드 수: {len(prs.slides)}")
    print(f"저장 위치: {output_path}")
    print("=" * 60)

if __name__ == "__main__":
    main()
